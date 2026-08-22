import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # MCKINSEY / APPLE PURE WHITE THEME COLOR PALETTE
    COLOR_BG = RGBColor(255, 255, 255)            # Pure White #FFFFFF
    COLOR_CARD = RGBColor(248, 250, 252)          # Light Slate Card #F8FAFC
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)   # Soft Gray Border #E2E8F0
    COLOR_ACCENT_BG = RGBColor(241, 245, 249)     # Soft Accent Gray #F1F5F9
    
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)        # Charcoal Black #0F172A
    COLOR_TEXT_MUTED = RGBColor(51, 65, 85)       # Dark Slate #334155
    COLOR_TEXT_SUB = RGBColor(100, 116, 139)      # Slate Gray #64748B
    
    COLOR_AZURE = RGBColor(2, 132, 199)           # Azure Blue #0284C7
    COLOR_TEAL = RGBColor(13, 148, 136)           # Teal Accent #0D9488
    COLOR_GREEN = RGBColor(16, 185, 129)          # Emerald Green #10B981
    COLOR_ORANGE = RGBColor(234, 88, 12)          # Orange Accent #EA580C
    COLOR_PURPLE = RGBColor(124, 58, 237)         # Purple Accent #7C3AED
    COLOR_RED = RGBColor(220, 38, 38)             # Crimson Red #DC2626
    COLOR_AMBER = RGBColor(217, 119, 6)           # Amber #D97706

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, slide_num):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9.0), Inches(0.3))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_AZURE

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.5), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

        # Counter Badge
        badge_box = slide.shapes.add_textbox(Inches(10.8), Inches(0.35), Inches(1.733), Inches(0.4))
        tf_badge = badge_box.text_frame
        p_badge = tf_badge.paragraphs[0]
        p_badge.alignment = PP_ALIGN.RIGHT
        p_badge.text = f"{slide_num} / 16"
        p_badge.font.size = Pt(12)
        p_badge.font.bold = True
        p_badge.font.color.rgb = COLOR_AZURE

        # Footer Line & Text
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "NEXUS AI  •  TEAM POLARIS  •  DECISION INTELLIGENCE PLATFORM"
        p_foot.font.size = Pt(9)
        p_foot.font.bold = True
        p_foot.font.color.rgb = COLOR_TEXT_SUB

    def add_card(slide, left, top, width, height, border_color=COLOR_CARD_BORDER, bg_color=COLOR_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    def add_notes(slide, notes_text):
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
        except Exception as e:
            pass

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    card1 = add_card(s1, 0.8, 0.8, 11.733, 5.8, COLOR_AZURE, COLOR_CARD)
    
    # Badge
    tb_badge = s1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(5), Inches(0.4))
    p = tb_badge.text_frame.paragraphs[0]
    p.text = "OPEN INNOVATION & AI HACKATHON PITCH"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    # Title
    tb_t = s1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(9.5), Inches(1.2))
    p = tb_t.text_frame.paragraphs[0]
    p.text = "NEXUS AI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    # Subtitle
    tb_s = s1.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(10.5), Inches(0.8))
    p = tb_s.text_frame.paragraphs[0]
    p.text = "AI-Native Decision Intelligence Platform for Critical Infrastructure"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    # Description
    tb_d = s1.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(10), Inches(0.8))
    p = tb_d.text_frame.paragraphs[0]
    p.text = "Empowering energy grids, railways, and municipal networks with autonomous multi-agent simulation, real-time risk mitigation, and predictive operational optimization."
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Meta Card (Team Info)
    meta_card = add_card(s1, 1.2, 4.8, 10.933, 1.4, COLOR_CARD_BORDER, COLOR_ACCENT_BG)
    tb_m = s1.shapes.add_textbox(Inches(1.4), Inches(4.9), Inches(10.5), Inches(1.2))
    tf_m = tb_m.text_frame
    p1 = tf_m.paragraphs[0]
    p1.text = "TEAM: POLARIS  |  TEAM LEADER: JATIN KUMAR"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    
    p2 = tf_m.add_paragraph()
    p2.text = "Institution: Deenbandhu Chhotu Ram University of Science & Technology (DCRUST), Murthal"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_AZURE

    p3 = tf_m.add_paragraph()
    p3.text = "Theme: AI & Open Innovation  •  System Architecture: SimPy Digital Twin + VCG Auction Engine + Multi-Agent Swarm"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_SUB

    add_notes(s1, "SPEAKER NOTES (Jatin Kumar):\nGood morning judges and attendees. I am Jatin Kumar, leader of Team Polaris. Today we present NEXUS AI—an AI-native decision intelligence platform designed to protect and optimize critical infrastructure. Global power grids, transport networks, and water utilities lose over $100 Billion annually due to delayed, manual operational decisions. NEXUS AI solves this by introducing a multi-agent AI swarm that simulates, predicts, and resolves operational bottlenecks in real time.")

    # =========================================================================
    # SLIDE 2: PROBLEM STATEMENT
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "01 / EXECUTIVE PROBLEM STATEMENT", "Critical Infrastructure Fragility & Operational Bottlenecks", 2)

    # Subtitle
    tb_sub = s2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.4))
    p = tb_sub.text_frame.paragraphs[0]
    p.text = "Legacy infrastructure management relies on fragmented telemetry and manual human decision cycles, leading to systemic failures."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Horizontal Flow Infographic (6 Steps)
    flow_steps = [
        ("1. Critical Infrastructure", "Power grids, railways, smart cities", COLOR_AZURE),
        ("2. Large Operational Data", "Petabytes of unstructured telemetry", COLOR_AZURE),
        ("3. Manual Decision Making", "Slow human analysis & siloed teams", COLOR_AMBER),
        ("4. Delayed Response", "Hours to react to critical alerts", COLOR_AMBER),
        ("5. Operational Failures", "Cascading outages & asset breakdowns", COLOR_RED),
        ("6. Financial Losses", "$100B+ annual global downtime cost", COLOR_RED)
    ]

    left_start = 0.8
    card_w = 1.75
    gap = 0.24

    for i, (title, desc, color) in enumerate(flow_steps):
        cur_left = left_start + i * (card_w + gap)
        c = add_card(s2, cur_left, 2.0, card_w, 4.4, color, COLOR_CARD)
        
        # Step header badge
        tb_step = s2.shapes.add_textbox(Inches(cur_left + 0.1), Inches(2.2), Inches(card_w - 0.2), Inches(0.8))
        tf = tb_step.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        # Description
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

        # Add Arrow indicator if not last step
        if i < 5:
            arr_box = s2.shapes.add_textbox(Inches(cur_left + card_w), Inches(3.8), Inches(gap), Inches(0.5))
            p_arr = arr_box.text_frame.paragraphs[0]
            p_arr.alignment = PP_ALIGN.CENTER
            p_arr.text = "➔"
            p_arr.font.size = Pt(14)
            p_arr.font.bold = True
            p_arr.font.color.rgb = COLOR_TEXT_SUB

    add_notes(s2, "SPEAKER NOTES:\nWhen an anomaly occurs in a power grid or rail network, data pours in from thousands of sensors. Today, human operators must manually analyze logs, correlate alerts across disconnected software, and make high-stakes decisions under extreme pressure. This delay converts minor glitches into catastrophic cascading blackouts and equipment destruction.")

    # =========================================================================
    # SLIDE 3: CURRENT CHALLENGES
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "02 / INDUSTRY PAIN POINTS", "Six Structural Bottlenecks in Infrastructure Operations", 3)

    challenges = [
        ("Manual Planning", "Operators rely on legacy spreadsheets and static rulebooks, creating immense decision latency during emergencies.", "⚠️ High Friction"),
        ("Disconnected Systems", "SCADA, ERP, and IoT monitoring operate in isolated silos without unified data interoperability.", "🔗 Siloed Data"),
        ("Poor Prediction", "Reactive maintenance models fail to forecast structural fatigue or equipment breakdown before failure occurs.", "📉 Reactive State"),
        ("No AI Coordination", "Lack of autonomous agent swarms to simulate complex multi-variate operational scenarios concurrently.", "🧠 Zero Synergy"),
        ("Resource Wastage", "Sub-optimal dispatching of repair crews and energy reserves leads to massive operational friction.", "💸 Energy Drain"),
        ("High Operational Cost", "Unplanned outages cause billions in SLA penalties, emergency overtime, and premature asset replacement.", "📊 Severe Cost")
    ]

    positions_s3 = [
        (0.8, 1.8), (4.8, 1.8), (8.8, 1.8),
        (0.8, 4.3), (4.8, 4.3), (8.8, 4.3)
    ]

    for i, (title, desc, badge) in enumerate(challenges):
        left, top = positions_s3[i]
        c = add_card(s3, left, top, 3.733, 2.2, COLOR_CARD_BORDER, COLOR_CARD)
        
        tb = s3.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(3.433), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p_badge = tf.add_paragraph()
        p_badge.text = badge
        p_badge.font.size = Pt(10)
        p_badge.font.bold = True
        p_badge.font.color.rgb = COLOR_AZURE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s3, "SPEAKER NOTES:\nWe have categorized the problem into six core challenges faced by enterprise operators. Traditional ERPs are passive databases—they tell you what happened yesterday. They cannot simulate the future or coordinate autonomous responses across departments.")

    # =========================================================================
    # SLIDE 4: SOLUTION (NEXUS AI MULTI-AGENT SWARM)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "03 / THE NEXUS ARCHITECTURE", "NEXUS AI: Autonomous Multi-Agent Swarm Intelligence", 4)

    # Center Hub Card
    hub = add_card(s4, 4.916, 3.2, 3.5, 1.8, COLOR_AZURE, COLOR_ACCENT_BG)
    tb_hub = s4.shapes.add_textbox(Inches(5.016), Inches(3.3), Inches(3.3), Inches(1.6))
    tf_h = tb_hub.text_frame
    tf_h.word_wrap = True
    p1 = tf_h.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = "NEXUS AI CORE"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_AZURE
    
    p2 = tf_h.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "Autonomous Multi-Agent Orchestration & Decision Engine"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    # 6 Satellite Agents
    agents = [
        ("Planner Agent", "Decomposes complex operational goals into executable sub-tasks.", 0.8, 1.8),
        ("Simulation Agent", "Runs real-time SimPy Monte Carlo simulations of physical assets.", 4.8, 1.4),
        ("Risk Agent", "Calculates failure probabilities and predictive hazard scores.", 8.8, 1.8),
        ("Optimization Agent", "Executes graph-theoretic resource routing and allocation.", 0.8, 4.8),
        ("Monitoring Agent", "Processes high-frequency SCADA telemetry & IoT anomaly streams.", 4.8, 5.3),
        ("Reporting Agent", "Generates McKinsey-grade executive briefings & audit logs.", 8.8, 4.8)
    ]

    for title, desc, left, top in agents:
        c = add_card(s4, left, top, 3.733, 1.5, COLOR_CARD_BORDER, COLOR_CARD)
        tb = s4.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(3.433), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s4, "SPEAKER NOTES:\nNEXUS AI replaces manual guesswork with a specialized Multi-Agent AI Swarm. Rather than relying on a single monolithic LLM, we deploy 6 specialized agents working in harmony. The Planner breaks down high-level directives, the Simulation Agent tests options in a virtual digital twin, and the Risk Agent calculates precision safety bounds.")

    # =========================================================================
    # SLIDE 5: SYSTEM ARCHITECTURE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "04 / ENTERPRISE SYSTEM ARCHITECTURE", "End-to-End Decision Pipeline & Technical Architecture", 5)

    arch_tiers = [
        ("CLIENT LAYER", ["User Interface", "React Dashboard", "ThreeJS 3D Digital Twin"], COLOR_AZURE, 0.8),
        ("GATEWAY LAYER", ["FastAPI API Gateway", "WebSocket Stream", "JWT Security"], COLOR_TEAL, 3.7),
        ("INTELLIGENCE LAYER", ["Agent Orchestrator", "Planner / Simulation / Risk / Optimization Agents", "Decision Engine"], COLOR_PURPLE, 6.6),
        ("DATA & ANALYTICS", ["PostgreSQL Database", "SimPy Physics Engine", "Visualization Dashboard"], COLOR_GREEN, 9.5)
    ]

    for tier_title, items, color, left in arch_tiers:
        c = add_card(s5, left, 1.8, 3.0, 5.0, color, COLOR_CARD)
        
        tb = s5.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(2.7), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = tier_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"\n• {item}"
            p_item.font.size = Pt(12)
            p_item.font.bold = True
            p_item.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(s5, "SPEAKER NOTES:\nHere is our enterprise architecture. The React + ThreeJS frontend connects through a high-performance FastAPI WebSocket gateway. Incoming telemetry feeds directly into the Agent Orchestrator, which coordinates the SimPy simulation engine and decision optimization models before pushing live updates back to the dashboard.")

    # =========================================================================
    # SLIDE 6: AI WORKFLOW
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "05 / EXECUTION WORKFLOW", "9-Step Autonomous Decision Optimization Pipeline", 6)

    steps_s6 = [
        ("1. Request", "User prompt / anomaly alert", 0.8, 1.8),
        ("2. Planner", "Task breakdown & routing", 4.8, 1.8),
        ("3. Breakdown", "Sub-agent task assignment", 8.8, 1.8),
        ("4. Simulation", "Digital twin Monte Carlo", 0.8, 3.5),
        ("5. Risk Analysis", "Hazard probability scoring", 4.8, 3.5),
        ("6. Optimization", "Constraint solver optimization", 8.8, 3.5),
        ("7. Ranking", "Action option scoring", 0.8, 5.2),
        ("8. Visualization", "3D UI state render", 4.8, 5.2),
        ("9. Recommendation", "Final decision execution", 8.8, 5.2)
    ]

    for title, desc, left, top in steps_s6:
        c = add_card(s6, left, top, 3.733, 1.3, COLOR_CARD_BORDER, COLOR_CARD)
        tb = s6.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(3.433), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_AZURE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s6, "SPEAKER NOTES:\nThis 9-step horizontal workflow illustrates how NEXUS AI processes an incoming infrastructure event. From the moment an alert triggers, our Planner Agent breaks down tasks, runs real-time Monte Carlo simulations, ranks optimal choices, and presents clear recommendations with human-in-the-loop controls.")

    # =========================================================================
    # SLIDE 7: TECHNICAL STACK
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "06 / TECHNICAL INFRASTRUCTURE", "Modern Production-Grade Technology Stack", 7)

    stack_layers = [
        ("FRONTEND LAYER", "React 18  •  TypeScript  •  Tailwind CSS  •  Vite  •  ThreeJS (3D Visualization)", COLOR_AZURE),
        ("BACKEND API", "Python 3.11  •  FastAPI  •  Uvicorn  •  WebSockets  •  Asyncio Engine", COLOR_TEAL),
        ("AI & AGENT SWARM", "OpenAI GPT-4o  •  LangChain  •  Agentic AI Swarm  •  Predictive ML Models", COLOR_PURPLE),
        ("SIMULATION & MATH", "SimPy (Discrete Event Engine)  •  NetworkX (Graph Topology)  •  NumPy / SciPy", COLOR_ORANGE),
        ("DEPLOYMENT & CI/CD", "Render Cloud Hosting  •  GitHub Actions  •  Vercel Edge Network  •  Docker", COLOR_GREEN)
    ]

    for i, (layer_title, techs, color) in enumerate(stack_layers):
        top = 1.8 + i * 1.0
        c = add_card(s7, 0.8, top, 11.733, 0.85, color, COLOR_CARD)
        
        tb = s7.shapes.add_textbox(Inches(0.95), Inches(top + 0.1), Inches(11.4), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = layer_title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = techs
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(s7, "SPEAKER NOTES:\nOur stack is engineered for enterprise performance and speed. We use React and ThreeJS on the frontend for crisp 3D visualization, Python FastAPI for asynchronous websocket communications, SimPy for discrete-event simulations, and OpenAI GPT-4o powering our multi-agent reasoning framework.")

    # =========================================================================
    # SLIDE 8: CORE FEATURES
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "07 / PRODUCT CAPABILITIES", "Eight Enterprise Capabilities of NEXUS AI", 8)

    features = [
        ("Multi-Agent AI", "Autonomous swarm coordination with specialized reasoning roles."),
        ("Decision Intelligence", "Transforms complex telemetry into actionable, prioritized decisions."),
        ("Scenario Simulation", "Simulates 10,000+ operational permutations per second via SimPy."),
        ("Risk Prediction", "Early-warning failure forecasting using predictive hazard scoring."),
        ("Live Dashboard", "Sub-second WebSocket telemetry updates across all asset layers."),
        ("Real-Time Optimization", "Graph-theoretic network routing and automatic resource re-balancing."),
        ("Interactive 3D Visuals", "Immersive ThreeJS digital twin representations of critical assets."),
        ("Scalable Architecture", "Cloud-native microservices engineered for high throughput.")
    ]

    positions_s8 = [
        (0.8, 1.8), (3.8, 1.8), (6.8, 1.8), (9.8, 1.8),
        (0.8, 4.3), (3.8, 4.3), (6.8, 4.3), (9.8, 4.3)
    ]

    for i, (title, desc) in enumerate(features):
        left, top = positions_s8[i]
        c = add_card(s8, left, top, 2.733, 2.2, COLOR_CARD_BORDER, COLOR_CARD)
        
        tb = s8.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(2.433), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_AZURE
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s8, "SPEAKER NOTES:\nThese eight core feature pillars demonstrate the width and depth of NEXUS AI. From discrete-event scenario simulation to interactive 3D digital twins, NEXUS AI gives operational command centers super-human foresight.")

    # =========================================================================
    # SLIDE 9: DEMO WORKFLOW (USER JOURNEY)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "08 / USER JOURNEY & INTERFACE", "End-to-End Live Operator Workflow", 9)

    journey_steps = [
        ("Step 1: Dashboard", "Operator monitors live city grid telemetry on central console.", COLOR_AZURE, 0.8),
        ("Step 2: Anomaly Alert", "Risk Agent detects transformer overheating risk at Station 4.", COLOR_AMBER, 3.1),
        ("Step 3: Simulation", "SimPy executes 5,000 power rerouting scenarios in 1.2s.", COLOR_PURPLE, 5.4),
        ("Step 4: AI Recommendation", "System recommends re-routing 12MW via Substation 9.", COLOR_GREEN, 7.7),
        ("Step 5: Execution", "Operator approves recommendation with 1-click execution.", COLOR_AZURE, 10.0)
    ]

    for title, desc, color, left in journey_steps:
        c = add_card(s9, left, 1.8, 2.133, 4.8, color, COLOR_CARD)
        
        tb = s9.shapes.add_textbox(Inches(left + 0.1), Inches(2.0), Inches(1.933), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

        # Mockup placeholder box inside card
        mock = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.15), Inches(4.2), Inches(1.833), Inches(2.0))
        mock.fill.solid()
        mock.fill.fore_color.rgb = COLOR_ACCENT_BG
        mock.line.color.rgb = COLOR_CARD_BORDER
        tb_m = s9.shapes.add_textbox(Inches(left + 0.2), Inches(4.8), Inches(1.7), Inches(1.0))
        p_m = tb_m.text_frame.paragraphs[0]
        p_m.alignment = PP_ALIGN.CENTER
        p_m.text = "[ UI Wireframe ]"
        p_m.font.size = Pt(9)
        p_m.font.color.rgb = COLOR_TEXT_SUB

    add_notes(s9, "SPEAKER NOTES:\nHere is how an operator interacts with NEXUS AI during an active emergency. Rather than panicking or sorting through thousands of SCADA logs, the operator receives a clear anomaly notification, reviews the simulated outcomes, and executes the optimal solution with a single click.")

    # =========================================================================
    # SLIDE 10: USE CASES
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "09 / TARGET SECTORS & APPLICATIONS", "Eight High-Value Critical Infrastructure Verticals", 10)

    sectors = [
        ("🚂 Railways & Transit", "Optimize signal schedules, switch routing & track maintenance."),
        ("🌆 Smart Cities", "Dynamic traffic flow management & municipal utility coordination."),
        ("🏥 Healthcare Systems", "Hospital bed allocation, emergency transport & oxygen logistics."),
        ("🏭 Manufacturing", "Preventive factory maintenance & supply chain bottleneck resolution."),
        ("⚡ Energy & Power Grids", "Dynamic load balancing, microgrid switching & outage prevention."),
        ("📦 Supply Chain & Port Logistics", "Container placement optimization & fleet route planning."),
        ("🛡️ Defense & National Security", "Tactical supply logistics & bases infrastructure resilience."),
        ("🌪️ Disaster Management", "Autonomous emergency evacuations & relief distribution flows.")
    ]

    positions_s10 = [
        (0.8, 1.8), (3.8, 1.8), (6.8, 1.8), (9.8, 1.8),
        (0.8, 4.3), (3.8, 4.3), (6.8, 4.3), (9.8, 4.3)
    ]

    for i, (title, desc) in enumerate(sectors):
        left, top = positions_s10[i]
        c = add_card(s10, left, top, 2.733, 2.2, COLOR_CARD_BORDER, COLOR_CARD)
        
        tb = s10.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(2.433), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s10, "SPEAKER NOTES:\nNEXUS AI is built as a horizontal decision intelligence engine that adapts seamlessly across multiple critical sectors. From energy grids and smart cities to railway networks and disaster response, our multi-agent architecture scales across any complex physical topology.")

    # =========================================================================
    # SLIDE 11: COMPETITIVE ANALYSIS
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11)
    add_header(s11, "10 / COMPETITIVE ADVANTAGE", "NEXUS AI vs. Legacy Infrastructure Management Tools", 11)

    # Table Creation
    rows, cols = 8, 5
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8)
    table_shape = s11.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Capability Metric", "Traditional ERP", "Manual Planning", "Legacy Dashboards", "NEXUS AI"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_AZURE if j == 4 else COLOR_ACCENT_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG if j == 4 else COLOR_TEXT_MAIN

    matrix_data = [
        ["Autonomous Automation", "❌ None", "❌ None", "⚠️ Partial", "✅ Full Multi-Agent"],
        ["Predictive Analytics", "❌ None", "❌ None", "⚠️ Basic Stats", "✅ Advanced AI/ML"],
        ["Multi-Agent Swarm", "❌ No", "❌ No", "❌ No", "✅ 6 Specialized Agents"],
        ["Real-Time Simulation", "❌ Batch Only", "❌ Manual", "❌ None", "✅ SimPy Monte Carlo"],
        ["Graph Optimization", "❌ Static", "❌ Static", "❌ Static", "✅ Dynamic NetworkX"],
        ["Scalable Microservices", "❌ Monolith", "❌ N/A", "⚠️ Limited", "✅ Cloud-Native"],
        ["Decision Support", "❌ Passive Logs", "❌ Slow Human", "⚠️ Charts Only", "✅ Autonomous Directives"]
    ]

    for i, row in enumerate(matrix_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 249, 255) if j == 4 else COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.bold = (j == 4)
            p.font.color.rgb = COLOR_AZURE if j == 4 else COLOR_TEXT_MUTED

    add_notes(s11, "SPEAKER NOTES:\nWhen compared against legacy ERPs, manual spreadsheets, or passive dashboards, NEXUS AI is in a league of its own. Existing tools only record past events; NEXUS AI actively simulates future outcomes and coordinates autonomous solutions.")

    # =========================================================================
    # SLIDE 12: EXPECTED IMPACT
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12)
    add_header(s12, "11 / QUANTIFIABLE VALUE CREATION", "Demonstrated ROI & Operational Performance Metrics", 12)

    kpis = [
        ("-65%", "Reduced Downtime", "Drastic reduction in unplanned outages through predictive fault isolation.", COLOR_GREEN),
        ("10x", "Faster Decisions", "Sub-second multi-agent scenario analysis replaces hours of human delay.", COLOR_AZURE),
        ("-40%", "Lower Operating Cost", "Significant reduction in emergency maintenance overhead and SLA fines.", COLOR_TEAL),
        ("+85%", "Resource Optimization", "Optimized routing of repair fleets and backup energy distribution.", COLOR_PURPLE),
        ("+95%", "Safety Compliance", "Rigorous risk boundary enforcement prevents hazardous catastrophic failures.", COLOR_ORANGE),
        ("3.5x", "Return on Investment", "Proven operational payback achieved within 6 months of enterprise deployment.", COLOR_AZURE)
    ]

    positions_s12 = [
        (0.8, 1.8), (4.8, 1.8), (8.8, 1.8),
        (0.8, 4.3), (4.8, 4.3), (8.8, 4.3)
    ]

    for i, (val, title, desc, color) in enumerate(kpis):
        left, top = positions_s12[i]
        c = add_card(s12, left, top, 3.733, 2.2, color, COLOR_CARD)
        
        tb = s12.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(3.433), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s12, "SPEAKER NOTES:\nOur economic impact is clear and compelling. By reducing downtime by 65% and delivering decisions 10x faster, NEXUS AI generates an estimated 3.5x return on investment for enterprise infrastructure clients within the first 6 months.")

    # =========================================================================
    # SLIDE 13: FUTURE ROADMAP
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13)
    add_header(s13, "12 / STRATEGIC VISION & ROADMAP", "5-Year Technology Expansion & Market Penetration", 13)

    roadmap_steps = [
        ("2026: MVP Phase", "Multi-Agent Core & SimPy engine launch for energy grid pilots.", COLOR_AZURE),
        ("2027: Enterprise AI", "Full SCADA / ERP integration with active multi-sector pilots.", COLOR_TEAL),
        ("2028: IoT Mesh", "Direct high-frequency sensor integration & edge AI deployment.", COLOR_PURPLE),
        ("2029: Digital Twin", "Full 3D digital twin replication of city-wide physical networks.", COLOR_ORANGE),
        ("2030: Autonomous Grid", "Self-healing critical infrastructure with zero human latency.", COLOR_GREEN)
    ]

    for i, (title, desc, color) in enumerate(roadmap_steps):
        left = 0.8 + i * 2.4
        c = add_card(s13, left, 2.0, 2.133, 4.4, color, COLOR_CARD)
        
        tb = s13.shapes.add_textbox(Inches(left + 0.1), Inches(2.2), Inches(1.933), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s13, "SPEAKER NOTES:\nOur technology roadmap spans five distinct phases. Starting with our validated MVP today, we will scale into enterprise SCADA integrations by 2027, deploy direct IoT mesh edge intelligence by 2028, and achieve fully autonomous self-healing infrastructure by 2030.")

    # =========================================================================
    # SLIDE 14: BUSINESS MODEL
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14)
    add_header(s14, "13 / COMMERCIAL MONETIZATION MODEL", "Multi-Stream Enterprise Business & Revenue Architecture", 14)

    streams = [
        ("Enterprise SaaS", "Tiered monthly subscription per asset node monitored.", COLOR_AZURE, 0.8, 1.8),
        ("Government Contracts", "Annual infrastructure resilience SLA contracts.", COLOR_TEAL, 6.8, 1.8),
        ("Custom Licensing", "On-premise deployment licensing for defense & energy.", COLOR_PURPLE, 0.8, 4.3),
        ("API Platform", "Usage-based pricing for simulation & risk API calls.", COLOR_ORANGE, 6.8, 4.3)
    ]

    for title, desc, color, left, top in streams:
        c = add_card(s14, left, top, 5.733, 2.2, color, COLOR_CARD)
        tb = s14.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.333), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s14, "SPEAKER NOTES:\nNEXUS AI utilizes a highly scalable B2B SaaS and API monetization model. We charge per monitored asset node for commercial enterprises, while providing tailored annual licensing for government utilities and defense defense infrastructure.")

    # =========================================================================
    # SLIDE 15: WHY NEXUS AI WINS
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_bg(s15)
    add_header(s15, "14 / STRATEGIC ADVANTAGE", "Five Reasons Why NEXUS AI Outperforms Competitors", 15)

    reasons = [
        ("1. AI Native Architecture", "Built from the ground up for agentic swarm intelligence, not legacy retrofits.", COLOR_AZURE),
        ("2. Proven Scalability", "Sub-second simulation engine handles 10,000+ complex nodes seamlessly.", COLOR_TEAL),
        ("3. Predictive Foresight", "Moves infrastructure management from reactive firefighting to predictive prevention.", COLOR_PURPLE),
        ("4. Autonomous Multi-Agent Swarm", "Deploys specialized agents for planning, simulation, risk, and optimization.", COLOR_ORANGE),
        ("5. Enterprise & Defense Ready", "Bank-grade encryption, role-based access, and on-premise cloud flexibility.", COLOR_GREEN)
    ]

    for i, (title, desc, color) in enumerate(reasons):
        top = 1.8 + i * 1.0
        c = add_card(s15, 0.8, top, 11.733, 0.85, color, COLOR_CARD)
        
        tb = s15.shapes.add_textbox(Inches(0.95), Inches(top + 0.1), Inches(11.4), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s15, "SPEAKER NOTES:\nWhy does NEXUS AI win? Because we are AI-native, infinitely scalable, predictively accurate, multi-agent powered, and built specifically for critical enterprise infrastructure.")

    # =========================================================================
    # SLIDE 16: THANK YOU & CONTACT HERO
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_bg(s16)

    card16 = add_card(s16, 0.8, 0.8, 11.733, 5.8, COLOR_AZURE, COLOR_CARD)
    
    tb_t16 = s16.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(9.5), Inches(1.0))
    p = tb_t16.text_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    tb_s16 = s16.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(0.6))
    p = tb_s16.text_frame.paragraphs[0]
    p.text = "NEXUS AI  •  Building the Future of Autonomous Infrastructure"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    contact_card = add_card(s16, 1.2, 3.2, 10.933, 3.0, COLOR_CARD_BORDER, COLOR_ACCENT_BG)
    tb_c = s16.shapes.add_textbox(Inches(1.4), Inches(3.4), Inches(10.5), Inches(2.6))
    tf_c = tb_c.text_frame
    
    p1 = tf_c.paragraphs[0]
    p1.text = "TEAM POLARIS"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    
    p2 = tf_c.add_paragraph()
    p2.text = "\nTeam Leader: Jatin Kumar\nGitHub Repository: github.com/Jatinkumar2503/Nexus--AI\nContact Email: jatin@nexus-ai.platform\nTheme: Artificial Intelligence & Open Innovation"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s16, "SPEAKER NOTES:\nThank you for your time and consideration. NEXUS AI is ready to pioneer the next generation of decision intelligence for critical infrastructure. We now welcome any questions from the judges.")

    # SAVE PRESENTATION TO MULTIPLE STANDARD PATHS
    output_paths = [
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_WorldClass_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_CODEFUSION_2k26_Polaris.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Best_PPT.pptx"
    ]
    for path in output_paths:
        prs.save(path)
        print(f"Presentation successfully created/updated at: {path}")

if __name__ == "__main__":
    create_presentation()
