"""NEXUS AI — Pure White Research Lab Style Visual Presentation Deck Generator.

Design Specification:
- Background: PURE WHITE (#FFFFFF) on EVERY slide
- Primary Text: DARK CHARCOAL (#202020)
- Secondary Text: DARK GRAY (#555555)
- Cards: VERY LIGHT GRAY (#F5F5F5)
- Borders: LIGHT GRAY (#D9D9D9)
- Primary Accent: SUBTLE DARK BROWN / GOLD (#8B6B23)
- Safe / Validated: MUTED GREEN (#2E7D32)
- Warning / Rejected: MUTED TERRACOTTA (#C62828)

Visual Structure:
Slide 1: Title & Railway Network Infographic (3 KPI cards, corridor visual)
Slide 2: Cascading Disruption vs NEXUS Recovery Flowcharts
Slide 3: End-to-End System Architecture (6 Layers, Data Flow Loop)
Slide 4: Counterfactual Simulation Flowchart & Policy Matrix Table
Slide 5: VCG Auction Diagram, Safety Verification Flowchart & Invariants Cards
Slide 6: Dispatcher Cockpit Mockup, CRLI Gauge Meter & 4-Tier Alert Funnel
Slide 7: Benchmark Charts (Delay Reduction, Regret, Latency Distribution, Severity Ladder)
Slide 8: 3 Visual Sections (Historical Replay, Property Safety Donut, Distillation Performance)
Slide 9: KPI Dashboard (8 Large Cards & End-to-End Latency Waterfall Chart)
Slide 10: 2x4 Innovation Matrix, Future Roadmap Flowchart & 7-Step Vision Chain
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# COLOR PALETTE (PURE WHITE RESEARCH LAB THEME)
COLOR_BG = RGBColor(255, 255, 255)         # PURE WHITE #FFFFFF
COLOR_CARD = RGBColor(245, 245, 245)       # VERY LIGHT GRAY #F5F5F5
COLOR_CARD_ALT = RGBColor(238, 238, 238)   # LIGHT GRAY CARD
COLOR_BORDER = RGBColor(217, 217, 217)     # LIGHT GRAY BORDER #D9D9D9
COLOR_TEXT_PRIMARY = RGBColor(32, 32, 32)  # DARK CHARCOAL #202020
COLOR_TEXT_MUTED = RGBColor(85, 85, 85)    # DARK GRAY #555555
COLOR_ACCENT = RGBColor(139, 107, 35)      # SUBTLE DARK BROWN / GOLD #8B6B23
COLOR_GREEN = RGBColor(46, 125, 50)        # MUTED GREEN #2E7D32
COLOR_RED = RGBColor(198, 40, 40)          # MUTED TERRACOTTA #C62828
COLOR_BLUE = RGBColor(21, 101, 192)        # DEEP CONTROL BLUE #1565C0
COLOR_PURPLE = RGBColor(106, 27, 154)      # DEEP PURPLE #6A1B9A

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT

    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRIMARY

def add_card(slide, left, top, width, height, border_color=COLOR_BORDER, bg_color=COLOR_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE & RAILWAY NETWORK INFOGRAPHIC
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Hero Pill Badge
    add_card(s1, 0.8, 0.6, 5.2, 0.4, COLOR_GREEN, COLOR_CARD_ALT)
    tb_pill = s1.shapes.add_textbox(Inches(0.9), Inches(0.65), Inches(5.0), Inches(0.3))
    p = tb_pill.text_frame.paragraphs[0]
    p.text = "🟢 LIVE AI PLATFORM  |  POLARIS DISPATCH ENGINE"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    # Title Box
    tb_t = s1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(7.0), Inches(2.2))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p1 = tf_t.paragraphs[0]
    p1.text = "NEXUS AI"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT

    p2 = tf_t.add_paragraph()
    p2.text = "Autonomous Decision Intelligence & VCG Dispatching for Rail Networks"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEXT_PRIMARY

    # 3 KPI Cards (Right Column)
    kpis = [
        ("318M / 1.45M", "Neural Model Scaling", COLOR_ACCENT),
        ("2.05 ms", "p50 Sub-5ms Latency", COLOR_GREEN),
        ("33.5%", "Historical Delay Reduction", COLOR_BLUE)
    ]
    for i, (val, label, col) in enumerate(kpis):
        top_pos = 1.2 + (i * 1.3)
        add_card(s1, 8.2, top_pos, 4.3, 1.15, COLOR_BORDER, COLOR_CARD)
        tb_kpi = s1.shapes.add_textbox(Inches(8.4), Inches(top_pos + 0.15), Inches(3.9), Inches(0.85))
        tf_k = tb_kpi.text_frame
        p_val = tf_k.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(26)
        p_val.font.bold = True
        p_val.font.color.rgb = col
        
        p_lbl = tf_k.add_paragraph()
        p_lbl.text = label
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED

    # Corridor Illustration Graphic (Bottom)
    add_card(s1, 0.8, 5.0, 11.73, 1.8, COLOR_BORDER, COLOR_CARD)
    tb_ill = s1.shapes.add_textbox(Inches(1.0), Inches(5.15), Inches(11.33), Inches(1.5))
    tf_i = tb_ill.text_frame
    tf_i.word_wrap = True
    p_i = tf_i.paragraphs[0]
    p_i.text = "🚆 HIGH-SPEED CORRIDOR INFOGRAPHIC MONITOR"
    p_i.font.size = Pt(12)
    p_i.font.bold = True
    p_i.font.color.rgb = COLOR_ACCENT

    points_ill = [
        "• Western Corridor Topology: Mumbai - Vadodara - Ahmedabad (500 km Dynamic Track Graph)",
        "• SimPy Discrete Event Simulation: Real-time telemetry ingestion at 1,000 steps/sec",
        "• Hard Safety Guarantee: Deterministic headway interlocks (≥120s) with 0% constraint violations"
    ]
    for pt in points_ill:
        p = tf_i.add_paragraph()
        p.text = pt
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 2: CASCADING DISRUPTION VS NEXUS RECOVERY FLOWCHARTS
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Problem vs Solution: Cascading Disruptions vs NEXUS AI Recovery", "Operational Challenge")

    # Left Flowchart Card: Baseline Cascading Disruption
    add_card(s2, 0.8, 1.6, 5.6, 5.3, COLOR_RED, COLOR_CARD)
    tb_fc1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_fc1 = tb_fc1.text_frame
    tf_fc1.word_wrap = True
    p = tf_fc1.paragraphs[0]
    p.text = "❌ TRADITIONAL CASCADING DISRUPTION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    steps_bad = [
        "Signal Failure (Virar Junction)",
        "↓",
        "Train A Delayed (+15 min)",
        "↓",
        "Platform 1 Occupancy Conflict",
        "↓",
        "Train B Delayed (+22 min)",
        "↓",
        "Train C Misses Scheduled Slot",
        "↓",
        "Network-Wide Gridlock (10+ Trains Affected)"
    ]
    for st in steps_bad:
        p = tf_fc1.add_paragraph()
        p.text = st
        p.font.size = Pt(10 if st == "↓" else 11)
        p.font.bold = True if st != "↓" else False
        p.font.color.rgb = COLOR_TEXT_MUTED if st == "↓" else COLOR_TEXT_PRIMARY

    # Right Flowchart Card: NEXUS Recovery
    add_card(s2, 6.9, 1.6, 5.6, 5.3, COLOR_GREEN, COLOR_CARD)
    tb_fc2 = s2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_fc2 = tb_fc2.text_frame
    tf_fc2.word_wrap = True
    p = tf_fc2.paragraphs[0]
    p.text = "🟢 NEXUS AI AUTONOMOUS RECOVERY"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    steps_good = [
        "Signal Failure Ingested (<5ms)",
        "↓",
        "Spatiotemporal Transformer Forward Pass",
        "↓",
        "Candidate Recovery Policies (A, B, C, D)",
        "↓",
        "Deterministic Safety Gate Verification",
        "↓",
        "Counterfactual Twin Evaluation (-34.5% Delay)",
        "↓",
        "Safe Detour Committed (Zero Domino Delays)"
    ]
    for st in steps_good:
        p = tf_fc2.add_paragraph()
        p.text = st
        p.font.size = Pt(10 if st == "↓" else 11)
        p.font.bold = True if st != "↓" else False
        p.font.color.rgb = COLOR_ACCENT if st == "↓" else COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 3: FULL SYSTEM ARCHITECTURE DIAGRAM
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "End-to-End NEXUS AI Closed-Loop Architecture", "System Design")

    layers = [
        ("1. Human Cockpit UI", "React + MapLibre GL + Attention Hub", COLOR_BLUE),
        ("2. FastAPI Gateway", "Async REST Endpoints & Streaming telemetry", COLOR_ACCENT),
        ("3. Digital Twin", "SimPy + NetworkX Physics Simulator", COLOR_PURPLE),
        ("4. Neural Core", "Spatiotemporal Transformer (1.45M / 318M)", COLOR_GREEN),
        ("5. Safety & VCG", "Deterministic Guardrails & VCG Token Auction", COLOR_ACCENT),
        ("6. Memory & Audit", "SQLite Recovery Memory & Replay Logger", COLOR_RED)
    ]

    for i, (name, desc, col) in enumerate(layers):
        left_pos = 0.8 + (i % 3) * 3.95
        top_pos = 1.8 + (i // 3) * 2.6
        add_card(s3, left_pos, top_pos, 3.7, 2.3, COLOR_BORDER, COLOR_CARD)
        
        tb = s3.shapes.add_textbox(Inches(left_pos + 0.15), Inches(top_pos + 0.15), Inches(3.4), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = name
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = col
        
        p1 = tf.add_paragraph()
        p1.text = f"\n{desc}"
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_TEXT_PRIMARY

    # Loop Banner (Bottom)
    add_card(s3, 0.8, 6.1, 11.73, 0.8, COLOR_BORDER, COLOR_CARD_ALT)
    tb_loop = s3.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(11.33), Inches(0.6))
    p = tb_loop.text_frame.paragraphs[0]
    p.text = "🔄 CLOSED-LOOP DATA FLOW: Human → Cockpit UI → FastAPI → Digital Twin → Neural Model → Safety Verification → Execution"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # =========================================================================
    # SLIDE 4: COUNTERFACTUAL SIMULATION & POLICY MATRIX TABLE
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Counterfactual Simulation & Candidate Policy Matrix", "Decision Engine")

    # Left Flowchart
    add_card(s4, 0.8, 1.6, 4.0, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_cf = s4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3.6), Inches(4.9))
    tf_cf = tb_cf.text_frame
    tf_cf.word_wrap = True
    p = tf_cf.paragraphs[0]
    p.text = "🔮 COUNTERFACTUAL FLOW"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    cf_steps = [
        "Current Railway State",
        "↓",
        "Generate Candidate Policies (A, B, C, D)",
        "↓",
        "SimPy Digital Twin Simulation",
        "↓",
        "Counterfactual Comparison",
        "↓",
        "Recommend Optimal Policy"
    ]
    for st in cf_steps:
        p = tf_cf.add_paragraph()
        p.text = st
        p.font.size = Pt(10 if st == "↓" else 11)
        p.font.bold = True if st != "↓" else False
        p.font.color.rgb = COLOR_ACCENT if st == "↓" else COLOR_TEXT_PRIMARY

    # Right Candidate Policy Table
    add_card(s4, 5.1, 1.6, 7.4, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_tbl = s4.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(7.0), Inches(4.9))
    tf_tbl = tb_tbl.text_frame
    tf_tbl.word_wrap = True
    p = tf_tbl.paragraphs[0]
    p.text = "📊 CANDIDATE RECOVERY POLICY COMPARISON MATRIX"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    policies = [
        ("Policy A (Hold 4 min)", "-10.5 min (-35.0%)", "2 Conflicts", "LOW Risk", "🟢 RECOMMENDED"),
        ("Policy B (Switch Loop)", "-7.5 min (-25.0%)", "4 Conflicts", "LOW Risk", "🟡 ADMISSIBLE"),
        ("Policy C (Speed Throttle)", "-4.5 min (-15.0%)", "6 Conflicts", "MED Risk", "🟡 ADMISSIBLE"),
        ("Policy D (Do Nothing)", "+6.0 min (+20.0%)", "10 Conflicts", "HIGH Risk", "🔴 REJECTED")
    ]

    for name, delay, conf, risk, status in policies:
        p = tf_tbl.add_paragraph()
        p.text = f"\n• {name}\n  Delay: {delay}  |  {conf}  |  {risk}  |  Status: {status}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN if "RECOMMENDED" in status else (COLOR_RED if "REJECTED" in status else COLOR_TEXT_PRIMARY)

    # =========================================================================
    # SLIDE 5: VCG AUCTION & SAFETY INVARIANTS CARDS
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Game-Theoretic VCG Auction & Formal Safety Invariants", "Safety & Allocation")

    # Left: VCG Auction Card
    add_card(s5, 0.8, 1.6, 5.6, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_vcg = s5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_v = tb_vcg.text_frame
    tf_v.word_wrap = True
    p = tf_v.paragraphs[0]
    p.text = "⚖️ VICKREY-CLARKE-GROVES (VCG) AUCTION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    vcg_pts = [
        "• Token Bidding Formula: b_i = Priority_i * (Delay_i + 1.0)",
        "• Vickrey Tax Rule: p_i* = ∑ V_j(a_-i*) - ∑ V_j(a*)",
        "• Truth-Telling Dominant Strategy: Prevents priority manipulation",
        "• Zero Freight Starvation: Cargo trains get fair slots without sidelining for 8+ hours"
    ]
    for pt in vcg_pts:
        p = tf_v.add_paragraph()
        p.text = f"\n{pt}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Right: Formal Invariants Card
    add_card(s5, 6.9, 1.6, 5.6, 5.3, COLOR_GREEN, COLOR_CARD)
    tb_inv = s5.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_inv = tb_inv.text_frame
    tf_inv.word_wrap = True
    p = tf_inv.paragraphs[0]
    p.text = "🛡️ FORMAL SAFETY INVARIANT CARDS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    invariants = [
        "1. Headway Invariant: Δt ≥ 120 seconds (Minimum Spatial Spacing)",
        "2. Braking Distance Invariant: d_brake = v² / (2a) ≤ d_clearance",
        "3. Route Exclusion Invariant: Route_A ∩ Route_B = ∅ (No Signal Conflicts)",
        "4. Platform Capacity Invariant: OccupiedPlatforms ≤ TotalCapacity"
    ]
    for inv in invariants:
        p = tf_inv.add_paragraph()
        p.text = f"\n{inv}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 6: DISPATCHER COCKPIT MOCKUP, CRLI GAUGE & TRIAGE FUNNEL
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Attention Management Hub: CRLI Meter & Alert Triage Funnel", "Human Factors")

    # Left: CRLI Meter Box
    add_card(s6, 0.8, 1.6, 5.6, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_crli = s6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_c = tb_crli.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "🧠 COGNITIVE REVIEW LOAD INDEX (CRLI 0–100)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    crli_info = [
        "CRLI Score = min(100, (10*D + 15*Complexity + 8*Q + 5*U) * W_mult)",
        "",
        "• 🟢 0 – 30 : QUIET STATE (Background Auto-Execution)",
        "• 🟡 35 – 74 : FOCUSED STATE (30s Batch Review Queues)",
        "• 🔴 75 – 100 : OVERLOAD STATE (Immediate Spotlight Popups)",
        "",
        "Smart Pre-fill Engine: Pre-fills hold times, platforms, and detour paths with 100% human dispatcher editability."
    ]
    for info in crli_info:
        p = tf_c.add_paragraph()
        p.text = info
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN if "0 – 30" in info else (COLOR_ACCENT if "35 – 74" in info else (COLOR_RED if "75 – 100" in info else COLOR_TEXT_PRIMARY))

    # Right: 4-Tier Alert Funnel Box
    add_card(s6, 6.9, 1.6, 5.6, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_fn = s6.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.9))
    tf_fn = tb_fn.text_frame
    tf_fn.word_wrap = True
    p = tf_fn.paragraphs[0]
    p.text = "🔻 ALERT PRIORITY TRIAGE FUNNEL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    funnel_steps = [
        "🟢 QUIET (Conf ≥ 85%, 0 Hazards) → Background Auto-Execute",
        "↓",
        "🟡 BATCH (Routine Advisories) → 30s Grouped Approvals",
        "↓",
        "🟠 PRIORITY (Medium Uncertainty) → Dispatcher Prompt",
        "↓",
        "🔴 IMMEDIATE (Safety Trip / OOD) → Spotlight Map Popups"
    ]
    for st in funnel_steps:
        p = tf_fn.add_paragraph()
        p.text = f"\n{st}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 7: REAL BENCHMARK CHARTS & SEVERITY LADDER
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "NEXUS-RailBench Scientific Evaluation & Latency Distribution", "Benchmark Results")

    # Native Bar Chart 1: Delay Reduction %
    cd1 = CategoryChartData()
    cd1.categories = ['FCFS', 'Fixed Priority', 'Greedy', 'CP-SAT MILP', 'NEXUS Edge']
    cd1.add_series('Delay Reduction %', (0.0, 8.5, 15.2, 34.5, 34.5))

    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.6)
    chart1 = s7.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, cd1).chart
    chart1.has_legend = False

    # Native Bar Chart 2: Latency Distribution (p50, p95, p99)
    cd2 = CategoryChartData()
    cd2.categories = ['p50 Latency', 'p95 Latency', 'p99 Latency']
    cd2.add_series('Latency (ms)', (3.62, 6.46, 7.02))

    x2, y2, cx2, cy2 = Inches(6.9), Inches(1.6), Inches(5.6), Inches(2.6)
    chart2 = s7.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x2, y2, cx2, cy2, cd2).chart
    chart2.has_legend = False

    # Severity Level Ladder Box (Bottom)
    add_card(s7, 0.8, 4.5, 11.73, 2.4, COLOR_BORDER, COLOR_CARD)
    tb_lad = s7.shapes.add_textbox(Inches(1.0), Inches(4.65), Inches(11.33), Inches(2.1))
    tf_l = tb_lad.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "🪜 NEXUS-RAILBENCH 5-LEVEL SEVERITY LADDER"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    ladder = [
        "Level 1 (Minor Delay - 4 min)  ➔  Level 2 (Moderate Rain - 12 min)  ➔  Level 3 (Heavy Rain - 25 min)",
        "Level 4 (Dense Fog - 45 min)   ➔  Level 5 (Network-Wide Gridlock - 90 min / 50 Trains)"
    ]
    for row in ladder:
        p = tf_l.add_paragraph()
        p.text = f"\n{row}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 8: 3 VISUAL SECTIONS (REPLAY, SAFETY DONUT, DISTILLATION)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Historical Replay, Safety Verification & Distillation Study", "Empirical Proof")

    # Section A: Historical Replay Card
    add_card(s8, 0.8, 1.6, 3.7, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_sa = s8.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(3.4), Inches(4.9))
    tf_sa = tb_sa.text_frame
    tf_sa.word_wrap = True
    p = tf_sa.paragraphs[0]
    p.text = "📜 SECTION A: HISTORICAL REPLAY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    pts_sa = [
        "• Incidents Replayed: 3 Major Indian Railways Disruption Events",
        "• Historical Human Delay: 574.0 min total",
        "• NEXUS Recovered Delay: 381.6 min total",
        "• Delay Savings: 33.5% Delay Reduction",
        "• Safety Constraint Trips: 0.0%"
    ]
    for pt in pts_sa:
        p = tf_sa.add_paragraph()
        p.text = f"\n{pt}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Section B: Property Safety Donut Chart Card
    add_card(s8, 4.8, 1.6, 3.7, 5.3, COLOR_GREEN, COLOR_CARD)
    tb_sb = s8.shapes.add_textbox(Inches(4.95), Inches(1.8), Inches(3.4), Inches(1.2))
    tf_sb = tb_sb.text_frame
    tf_sb.word_wrap = True
    p = tf_sb.paragraphs[0]
    p.text = "🛡️ SECTION B: PROPERTY SAFETY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    # Native Doughnut Chart for Safety States
    cd_donut = CategoryChartData()
    cd_donut.categories = ['Safe Accepted', 'Unsafe Blocked']
    cd_donut.add_series('States', (41981, 58019))

    x_d, y_d, cx_d, cy_d = Inches(4.95), Inches(3.0), Inches(3.4), Inches(3.6)
    chart_d = s8.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x_d, y_d, cx_d, cy_d, cd_donut).chart
    chart_d.has_legend = True
    chart_d.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Section C: Knowledge Distillation Card
    add_card(s8, 8.8, 1.6, 3.7, 5.3, COLOR_BORDER, COLOR_CARD)
    tb_sc = s8.shapes.add_textbox(Inches(8.95), Inches(1.8), Inches(3.4), Inches(4.9))
    tf_sc = tb_sc.text_frame
    tf_sc.word_wrap = True
    p = tf_sc.paragraphs[0]
    p.text = "🧠 SECTION C: DISTILLATION STUDY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    pts_sc = [
        "Teacher Model (318M)  ➔  Student Model (1.45M)",
        "• MAE Improvement: 0.460m  ➔  0.220m",
        "• Policy Regret Drop: 0.12  ➔  0.04",
        "• Latency Preserved: 2.10 ms p50",
        "• Memory Footprint: 5.6 MB"
    ]
    for pt in pts_sc:
        p = tf_sc.add_paragraph()
        p.text = f"\n{pt}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 9: KPI DASHBOARD & END-TO-END LATENCY WATERFALL CHART
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "KPI Executive Dashboard & End-to-End Latency Waterfall", "System Performance")

    # 8 Large KPI Cards (2x4 Grid on Left)
    kpi_cards = [
        ("318M", "Heavy Model Params", COLOR_ACCENT),
        ("1.45M", "Edge Model Params", COLOR_GREEN),
        ("5.6 MB", "Edge Checkpoint", COLOR_PURPLE),
        ("2.10 ms", "Sub-3ms p50 Speed", COLOR_BLUE),
        ("3.297 ms", "End-to-End Latency", COLOR_ACCENT),
        ("7.937 ms", "p95 Tail Latency", COLOR_BLUE),
        ("33.5%", "Delay Reduction", COLOR_GREEN),
        ("0 Violations", "Formal Invariants", COLOR_RED)
    ]

    for i, (val, label, col) in enumerate(kpi_cards):
        left_pos = 0.8 + (i % 4) * 1.4
        top_pos = 1.6 + (i // 4) * 1.3
        add_card(s9, left_pos, top_pos, 1.3, 1.15, COLOR_BORDER, COLOR_CARD)
        tb_k = s9.shapes.add_textbox(Inches(left_pos + 0.05), Inches(top_pos + 0.1), Inches(1.2), Inches(0.95))
        tf_k = tb_k.text_frame
        p_v = tf_k.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(14)
        p_v.font.bold = True
        p_v.font.color.rgb = col
        
        p_l = tf_k.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(8)
        p_l.font.color.rgb = COLOR_TEXT_MUTED

    # Latency Waterfall Chart (Right)
    cd_water = CategoryChartData()
    cd_water.categories = ['Inference', 'Feature Prep', 'Policy Gen', 'Safety Verify', 'REST API', 'Total E2E']
    cd_water.add_series('Latency (ms)', (1.25, 0.45, 0.65, 0.55, 0.723, 3.297))

    xw, yw, cxw, cyw = Inches(6.6), Inches(1.6), Inches(5.9), Inches(5.3)
    chart_w = s9.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, xw, yw, cxw, cyw, cd_water).chart
    chart_w.has_legend = False

    # =========================================================================
    # SLIDE 10: INNOVATION MATRIX, ROADMAP FLOWCHART & 7-STEP VISION CHAIN
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Innovation Matrix, Future Roadmap & 7-Step Vision Chain", "Future Vision")

    # 2x4 Innovation Matrix Card
    add_card(s10, 0.8, 1.6, 5.6, 2.5, COLOR_BORDER, COLOR_CARD)
    tb_im = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(2.3))
    tf_im = tb_im.text_frame
    tf_im.word_wrap = True
    p = tf_im.paragraphs[0]
    p.text = "⚡ 2x4 INNOVATION MATRIX"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    inn_pts = [
        "1. Spatiotemporal GAT  |  2. TorchScript JIT  |  3. VCG Token Auction  |  4. CRLI Load Index",
        "5. Distillation Engine  |  6. Counterfactual Twin  |  7. Invariants Safety  |  8. RailBench Suite"
    ]
    for pt in inn_pts:
        p = tf_im.add_paragraph()
        p.text = f"\n{pt}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Future Roadmap Flowchart Card
    add_card(s10, 6.9, 1.6, 5.6, 2.5, COLOR_GREEN, COLOR_CARD)
    tb_rm = s10.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(2.3))
    tf_rm = tb_rm.text_frame
    tf_rm.word_wrap = True
    p = tf_rm.paragraphs[0]
    p.text = "🗺️ FUTURE ROADMAP FLOWCHART"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    rm_steps = [
        "Real Railway Data Ingestion  ➔  Independent Validation",
        "➔  Multi-Agent Coordination  ➔  Confidence-Aware Autonomy",
        "➔  Conditional Autonomous Dispatch Deployment"
    ]
    for pt in rm_steps:
        p = tf_rm.add_paragraph()
        p.text = f"\n{pt}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_PRIMARY

    # Final 7-Step Vision Chain (Bottom)
    add_card(s10, 0.8, 4.3, 11.73, 2.6, COLOR_BORDER, COLOR_CARD_ALT)
    tb_v = s10.shapes.add_textbox(Inches(1.0), Inches(4.45), Inches(11.33), Inches(2.3))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    p = tf_v.paragraphs[0]
    p.text = "🎯 FINAL 7-STEP DECISION VISION CHAIN"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    chain_steps = [
        "PREDICT  ➔  SIMULATE  ➔  OPTIMIZE  ➔  VERIFY  ➔  EXPLAIN  ➔  ASSIST  ➔  LEARN",
        "",
        "🚀 TEAM POLARIS  |  KING LEADER: JATIN KUMAR  |  DCRUST MURTHAL",
        "🌐 GitHub Repository: https://github.com/Jatinkumar2503/Nexus--AI"
    ]
    for st in chain_steps:
        p = tf_v.add_paragraph()
        p.text = st
        p.font.size = Pt(12 if "PREDICT" in st else (14 if "TEAM POLARIS" in st else 10))
        p.font.bold = True
        p.font.color.rgb = COLOR_GREEN if "PREDICT" in st else (COLOR_ACCENT if "TEAM POLARIS" in st else COLOR_TEXT_MUTED)

    # SAVE TO PERMISSION SAFE PATHS
    paths = [
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Ultimate_Visual_Deck.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Futuristic_Deck.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_10_Slides_Futuristic.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Best_PPT.pptx"
    ]
    for pth in paths:
        try:
            prs.save(pth)
            print(f"Pure White Research Deck saved to: {pth}")
        except PermissionError:
            print(f"Skipped open file: {pth}")
        except Exception as e:
            print(f"Error saving to {pth}: {e}")

if __name__ == "__main__":
    create_deck()
