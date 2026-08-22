import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_open_layout_10_slide_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # MODERN ELEGANT COLOR PALETTE (CLEAN, BOX-FREE ARCHITECTURE)
    COLOR_BG = RGBColor(255, 255, 255)            # Pure White Background
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)        # Deep Charcoal Black #0F172A
    COLOR_TEXT_MUTED = RGBColor(51, 65, 85)       # Slate Gray #334155
    COLOR_TEXT_SUB = RGBColor(100, 116, 139)      # Soft Slate #64748B
    
    # ELEGANT ACCENT COLORS
    COLOR_AZURE = RGBColor(2, 132, 199)           # Primary Azure Blue #0284C7
    COLOR_TEAL = RGBColor(13, 148, 136)           # Teal Accent #0D9488
    COLOR_PURPLE = RGBColor(124, 58, 237)         # Deep Purple #7C3AED
    COLOR_EMERALD = RGBColor(16, 185, 129)        # Emerald Green #10B981
    COLOR_RED = RGBColor(220, 38, 38)             # Crimson Red #DC2626
    COLOR_AMBER = RGBColor(217, 119, 6)           # Warm Amber #D97706
    COLOR_BORDER_LINE = RGBColor(226, 232, 240)   # Light Gray Rule #E2E8F0

    def set_clean_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, slide_num):
        # Category Tag (Open Text)
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.0), Inches(0.3))
        tf_tag = tag_box.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_AZURE

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10.5), Inches(0.6))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

        # Slide Counter Number (Open Right Alignment)
        counter_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.4), Inches(1.533), Inches(0.4))
        tf_c = counter_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.RIGHT
        p_c.text = f"{slide_num:02d} / 10"
        p_c.font.size = Pt(12)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_TEXT_SUB

        # Subtle Horizontal Divider Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER_LINE
        line.line.fill.background()

        # Footer Text
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "NEXUS AI  |  TEAM POLARIS  |  DECISION INTELLIGENCE PLATFORM"
        p_foot.font.size = Pt(9)
        p_foot.font.bold = True
        p_foot.font.color.rgb = COLOR_TEXT_SUB

    def add_notes(slide, notes_text):
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
        except Exception as e:
            pass

    # =========================================================================
    # SLIDE 1: OPEN HERO TITLE SLIDE (NO BOXES, PURE TYPOGRAPHY & WHITESPACE)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s1)

    # Category Badge Text
    tb_pill = s1.shapes.add_textbox(Inches(1.3), Inches(1.15), Inches(5.0), Inches(0.3))
    p = tb_pill.text_frame.paragraphs[0]
    p.text = "NEXUS AI PLATFORM  •  OPEN INNOVATION PITCH"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    # Huge Main Title
    tb_t = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(1.2))
    p = tb_t.text_frame.paragraphs[0]
    p.text = "NEXUS AI"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    # Subtitle
    tb_s = s1.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.8))
    p = tb_s.text_frame.paragraphs[0]
    p.text = "AI-Native Decision Intelligence & VCG Token Auction Engine"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    # Description Paragraph
    tb_d = s1.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.9))
    p = tb_d.text_frame.paragraphs[0]
    p.text = "Empowering railway dispatchers to predict, simulate, and resolve critical network disruptions in real time using a SimPy digital twin and a multi-agent AI swarm."
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Elegant Horizontal Accent Divider
    div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.733), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_AZURE
    div.line.fill.background()

    # Open Team Details Block (No Bounding Box)
    tb_team = s1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.733), Inches(1.8))
    tf_team = tb_team.text_frame
    
    p1 = tf_team.paragraphs[0]
    p1.text = "TEAM POLARIS   |   LEADER: JATIN KUMAR"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf_team.add_paragraph()
    p2.text = "Deenbandhu Chhotu Ram University of Science and Technology (DCRUST), Murthal"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_AZURE

    p3 = tf_team.add_paragraph()
    p3.text = "GitHub Repository: github.com/Jatinkumar2503/Nexus--AI   •   Theme: Artificial Intelligence & Infrastructure Resilience"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_SUB

    add_notes(s1, "SPEAKER NOTES (Jatin Kumar):\nGood morning judges and attendees. I am Jatin Kumar, leader of Team Polaris from DCRUST Murthal. Today we present NEXUS AI—an AI-native decision intelligence platform designed to protect and optimize critical railway infrastructure. Global transport networks lose over $100 Billion annually due to delayed, manual operational guesswork during unexpected track disruptions.")

    # =========================================================================
    # SLIDE 2: THE PROBLEM (OPEN TIMELINE / NUMBERED FLOW, NO BOXES)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s2)
    add_header(s2, "01 / EXECUTIVE PROBLEM STATEMENT", "The Cascading Cost of Railway Disruptions", 2)

    # Subtitle
    tb_sub = s2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.4))
    p = tb_sub.text_frame.paragraphs[0]
    p.text = "Fragmented operational telemetry and manual dispatcher guesswork create cascading network failures."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    steps = [
        ("01", "Signal Glitch", "Track circuit or signal failure at a critical junction node.", COLOR_AZURE),
        ("02", "Telemetry Overload", "Petabytes of raw sensor data overwhelm human dispatchers.", COLOR_AZURE),
        ("03", "Manual Guesswork", "Slow phone calls and manual route recalculation under pressure.", COLOR_AMBER),
        ("04", "Crew & Safety Breach", "Train crews exceed legal shift limits; gridlock propagates.", COLOR_RED),
        ("05", "$100B+ Global Damage", "Cascading delays burn traction power & disrupt passenger schedules.", COLOR_RED)
    ]

    col_w = 2.15
    gap = 0.24

    for i, (num, title, desc, color) in enumerate(steps):
        c_left = 0.8 + i * (col_w + gap)
        
        # Giant Step Number
        tb_num = s2.shapes.add_textbox(Inches(c_left), Inches(2.1), Inches(col_w), Inches(0.6))
        p = tb_num.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

        # Step Title
        tb_t = s2.shapes.add_textbox(Inches(c_left), Inches(2.8), Inches(col_w), Inches(0.6))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        p = tf_t.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        # Top Stripe Indicator
        stripe = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(c_left), Inches(3.5), Inches(1.8), Inches(0.04))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()

        # Step Description Body
        tb_d = s2.shapes.add_textbox(Inches(c_left), Inches(3.7), Inches(col_w), Inches(3.0))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p = tf_d.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s2, "SPEAKER NOTES:\nWhen a signal fails or a track is blocked, human dispatchers must manually analyze complex train schedules under intense pressure. This causes delayed reactions, resulting in cascading train standstills, crew overtime violations, and massive energy waste.")

    # =========================================================================
    # SLIDE 3: THE SOLUTION OVERVIEW (OPEN 3-ROW ARCHITECTURE, NO ENCLOSED CARDS)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s3)
    add_header(s3, "02 / SOLUTION OVERVIEW", "Predict, Simulate & Decide Before Impact", 3)

    pillars = [
        ("SimPy Physics Digital Twin", "Real-time discrete-event engine modeling exact train kinematics, platform track capacities, and IEEE catenary power limits.", COLOR_AZURE),
        ("VCG Token Slot Auction Engine", "Game-theoretic priority token bidding allocating track slots fairly with second-price pricing and anti-starvation mechanics.", COLOR_TEAL),
        ("LangGraph Multi-Agent AI Swarm", "Cyclic state machine orchestrating Planner, Risk, Energy, and Safety Validation agents to negotiate Pareto-optimal recovery plans.", COLOR_PURPLE)
    ]

    for i, (title, desc, color) in enumerate(pillars):
        top = 1.7 + i * 1.65

        # Left Vertical Color Accent Line
        acc = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.08), Inches(1.2))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        # Title & Content
        tb = s3.shapes.add_textbox(Inches(1.1), Inches(top - 0.05), Inches(11.4), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s3, "SPEAKER NOTES:\nNEXUS AI introduces three core breakthroughs: A real-time SimPy Digital Twin that models network physics, a VCG Token Auction Engine that allocates track priority mathematically, and a LangGraph Multi-Agent Swarm that generates validated recovery scenarios.")

    # =========================================================================
    # SLIDE 4: SYSTEM PIPELINE (OPEN HORIZONTAL FLOW)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s4)
    add_header(s4, "03 / SYSTEM ARCHITECTURE", "End-to-End Edge Telemetry to Cockpit Decision Flow", 4)

    pipeline_steps = [
        ("01", "Edge Sensors & Azure IoT", "Ingests live track telemetry, train velocity, and disruption events.", COLOR_AZURE),
        ("02", "SimPy Digital Twin", "Simulates network movement, platform capacity, & catenary power load.", COLOR_TEAL),
        ("03", "LangGraph AI Swarm", "Multi-agent negotiation formulating Pareto recovery options.", COLOR_PURPLE),
        ("04", "VCG Auction Solver", "Executes second-price slot bidding & anti-starvation rules.", COLOR_AMBER),
        ("05", "Dispatcher Cockpit", "Presents validated recommendations for 1-click human approval.", COLOR_EMERALD)
    ]

    for i, (num, title, desc, color) in enumerate(pipeline_steps):
        top = 1.6 + i * 1.05

        # Horizontal accent indicator line
        acc = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top + 0.05), Inches(0.08), Inches(0.75))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        tb = s4.shapes.add_textbox(Inches(1.1), Inches(top), Inches(11.4), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"{num}.  {title}"
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s4, "SPEAKER NOTES:\nTelemetry flows from edge track sensors to our SimPy Digital Twin. The LangGraph Multi-Agent Swarm negotiates recovery options, the VCG Solver resolves track slot conflicts, and the optimal strategy is presented to the human dispatcher.")

    # =========================================================================
    # SLIDE 5: VCG TOKEN AUCTION ENGINE (OPEN 3-COLUMN TYPOGRAPHY)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s5)
    add_header(s5, "04 / GAME THEORY ENGINE", "Vickrey-Clarke-Groves (VCG) Token Slot Auctions", 5)

    vcg_cols = [
        ("01", "Priority Tokens (🎫 tkn)", "Train agents submit priority tokens reflecting delay penalties. High-priority express trains bid higher for critical platform track access.", COLOR_AZURE),
        ("02", "Second-Price Pricing Rule", "The highest bidder wins the track slot, but pays only the second-highest bid. Guarantees truthful bidding where no train benefits by over-reporting.", COLOR_TEAL),
        ("03", "Anti-Starvation Mechanism", "Trains waiting at red signals continuously accrue token interest over time, ensuring low-priority freight trains eventually win slots.", COLOR_AMBER)
    ]

    col_w = 3.6
    gap = 0.46

    for i, (num, title, desc, color) in enumerate(vcg_cols):
        c_left = 0.8 + i * (col_w + gap)

        # Big Number Header
        tb_n = s5.shapes.add_textbox(Inches(c_left), Inches(1.7), Inches(col_w), Inches(0.5))
        p = tb_n.text_frame.paragraphs[0]
        p.text = f"MECHANISM {num}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        # Accent Bar
        bar = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(c_left), Inches(2.2), Inches(3.2), Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Title
        tb_t = s5.shapes.add_textbox(Inches(c_left), Inches(2.35), Inches(col_w), Inches(0.8))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        p = tf_t.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

        # Description Text
        tb_d = s5.shapes.add_textbox(Inches(c_left), Inches(3.2), Inches(col_w), Inches(3.5))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p = tf_d.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s5, "SPEAKER NOTES:\nTo solve platform track conflicts fairly, NEXUS AI employs VCG token auctions. Train agents submit priority tokens based on delay penalties. Second-Price bidding encourages truthful reporting, while interest token accumulation prevents lower-priority trains from starving.")

    # =========================================================================
    # SLIDE 6: MULTI-AGENT ORCHESTRATION SWARM (OPEN 2x2 QUADRANT)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s6)
    add_header(s6, "05 / MULTI-AGENT INTELLIGENCE", "LangGraph Cyclic Multi-Agent Negotiation Swarm", 6)

    agents = [
        ("Planner Agent", "Formulates multi-path recovery hypotheses (Detour, Short-Turn, Holds).", COLOR_AZURE, 0.8, 1.8),
        ("Risk & Energy Agent", "Calculates traction acceleration power, fuel burn & IEEE catenary limits.", COLOR_TEAL, 6.8, 1.8),
        ("VCG Auction Solver", "Executes game-theoretic second-price priority token bidding.", COLOR_AMBER, 0.8, 4.3),
        ("Validation & Safety Agent", "Enforces 100% deterministic safety rules & self-reflection bounds.", COLOR_PURPLE, 6.8, 4.3)
    ]

    for title, desc, color, left, top in agents:
        # Accent Stripe Left
        acc = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(1.8))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        tb = s6.shapes.add_textbox(Inches(left + 0.2), Inches(top), Inches(5.5), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s6, "SPEAKER NOTES:\nOur AI swarm consists of specialized agents working in a cyclic LangGraph state machine. The Planner proposes routes, Risk & Energy verify physical limits, VCG Auction resolves slot order, and Validation guarantees safety before dispatcher review.")

    # =========================================================================
    # SLIDE 7: SCENARIO COMPARISON MATRIX (OPEN TABLE LAYOUT, NO ROUNDED CARDS)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s7)
    add_header(s7, "06 / DECISION INTELLIGENCE", "Multi-Objective Pareto-Optimal Tradeoff Analysis", 7)

    scenarios = [
        ("Option A: Do Nothing", "HIGH RISK / REJECTED", "140+ Min Delay  •  Severe Crew Overtime Violation  •  Cascading Station Gridlock", COLOR_RED),
        ("Option B: Detour Route", "PARETO OPTIMAL / RECOMMENDED", "22 Min Delay  •  0 Min Station Gridlock  •  100% Catenary Power Safety  •  Optimal Pareto Score", COLOR_EMERALD),
        ("Option C: Short-Turn", "FEASIBLE ALTERNATIVE", "15 Min Delay  •  Frees Platform Instantly  •  Requires Passenger Connection Transfer", COLOR_AMBER)
    ]

    for i, (title, badge, desc, color) in enumerate(scenarios):
        top = 1.8 + i * 1.65

        # Left stripe bar
        acc = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.08), Inches(1.3))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        tb = s7.shapes.add_textbox(Inches(1.1), Inches(top - 0.05), Inches(11.4), Inches(1.35))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"{title}   —   [{badge}]"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s7, "SPEAKER NOTES:\nNEXUS AI doesn't just offer one solution—it provides side-by-side Pareto-optimal comparisons. Here, Option B (Detour Route) is recommended because it completely eliminates station gridlock while maintaining full safety scores.")

    # =========================================================================
    # SLIDE 8: DISPATCHER COCKPIT & SAFEGUARDS (OPEN LIST WITH TYPOGRAPHY STACK)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s8)
    add_header(s8, "07 / OPERATOR EXPERIENCE", "Human-in-the-Loop Cockpit & Durable Audit Memory", 8)

    features = [
        ("Interactive MapLibre GL Corridor Map", "Live visual rendering of train positions, track signals, and dynamic bottleneck locations.", COLOR_AZURE),
        ("Plain-Language AI Reasoning", "Explains exactly *why* a plan was selected with side-by-side scenario tradeoff analysis.", COLOR_TEAL),
        ("1-Click Dispatcher Approval Gate", "Simulation resumes only after the human operator explicitly reviews and approves the plan.", COLOR_EMERALD),
        ("Durable Event Replay Memory", "Immutable audit log tracking every incident telemetry event, token bid, and executed recovery action.", COLOR_PURPLE)
    ]

    for i, (title, desc, color) in enumerate(features):
        top = 1.7 + i * 1.25

        acc = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top + 0.05), Inches(0.08), Inches(0.8))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()

        tb = s8.shapes.add_textbox(Inches(1.1), Inches(top), Inches(11.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s8, "SPEAKER NOTES:\nThe React dispatcher cockpit keeps human operators in complete control. It displays spatial train positions, plain-language reasoning, requires explicit dispatcher approval, and maintains an immutable audit trail for post-incident review.")

    # =========================================================================
    # SLIDE 9: EMPIRICAL BENCHMARKS & TECH STACK (MASSIVE HERO NUMBERS, NO BOXES)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s9)
    add_header(s9, "08 / BENCHMARKS & TECH STACK", "Validated Metrics & Production Architecture", 9)

    # Left Column: Huge Metrics Typography
    metrics = [
        ("78%", "Cascading Delay Reduction", "Cascading network delay reduced from 140m to <30m.", COLOR_AZURE),
        ("4.2x", "Faster Incident Recovery", "Resolution time reduced from 45m to <11m.", COLOR_TEAL),
        ("99.98%", "Rule Safety Compliance", "Zero physical collision or power over-subscription errors.", COLOR_EMERALD)
    ]

    for i, (num, label, desc, color) in enumerate(metrics):
        top = 1.7 + i * 1.65

        # Massive Number
        tb_n = s9.shapes.add_textbox(Inches(0.8), Inches(top), Inches(2.2), Inches(1.2))
        p = tb_n.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = color

        # Label & Text
        tb_t = s9.shapes.add_textbox(Inches(3.0), Inches(top + 0.1), Inches(3.6), Inches(1.2))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True

        p1 = tf_t.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN

        p2 = tf_t.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Vertical Middle Divider Line
    div = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.7), Inches(0.02), Inches(4.8))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_BORDER_LINE
    div.line.fill.background()

    # Right Column: Production Tech Stack List
    tb_st = s9.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.3), Inches(4.8))
    tf_st = tb_st.text_frame
    tf_st.word_wrap = True

    p = tf_st.paragraphs[0]
    p.text = "PRODUCTION TECH STACK"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    stack = [
        ("Frontend UI", "React 18, TypeScript, MapLibre GL, Tailwind CSS"),
        ("Backend Services", "Python 3.11+, FastAPI, WebSockets, SimPy, NetworkX"),
        ("AI Core & Solvers", "LangGraph, OpenAI GPT-4o, VCG Auction Engine"),
        ("Cloud Infrastructure", "Docker, Azure IoT Hub, Render Cloud Infrastructure")
    ]
    for title, desc in stack:
        p1 = tf_st.add_paragraph()
        p1.text = f"\n• {title}"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN

        p2 = tf_st.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s9, "SPEAKER NOTES:\nEmpirical benchmarks demonstrate a 78% reduction in network delays and 4.2x faster incident recovery. Built using standard enterprise technologies: React, FastAPI, SimPy, NetworkX, and LangGraph.")

    # =========================================================================
    # SLIDE 10: STRATEGIC ADVANTAGE & THANK YOU (OPEN HERO CONCLUSION)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_clean_bg(s10)
    add_header(s10, "09 / SUMMARY & CONCLUSION", "Transforming Critical Infrastructure Operations", 10)

    # Title
    tb_t = s10.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(0.5))
    p = tb_t.text_frame.paragraphs[0]
    p.text = "WHY NEXUS AI WINS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZURE

    points = [
        "1. AI-Native Swarm Architecture: Built specifically for autonomous multi-agent infrastructure orchestration.",
        "2. Mathematically Fair VCG Auctions: Zero starvation, truthful slot allocation for competing train priorities.",
        "3. Deterministic Safety Safeguards: Local fallback engine guarantees 100% operational compliance.",
        "4. Proven Scale & Impact: 78% delay reduction with sub-second discrete-event simulation performance."
    ]

    tb_p = s10.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.733), Inches(2.4))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True

    for i, pt in enumerate(points):
        p_pt = tf_p.paragraphs[0] if i == 0 else tf_p.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = COLOR_TEXT_MUTED

    # Horizontal Divider Line
    div_c = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.733), Inches(0.03))
    div_c.fill.solid()
    div_c.fill.fore_color.rgb = COLOR_AZURE
    div_c.line.fill.background()

    # Open Contact Text
    tb_c = s10.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.6))
    tf_c = tb_c.text_frame

    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "THANK YOU   |   TEAM POLARIS   |   LEADER: JATIN KUMAR"
    p_c1.font.size = Pt(18)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_TEXT_MAIN

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Deenbandhu Chhotu Ram University of Science & Technology (DCRUST), Murthal"
    p_c2.font.size = Pt(13)
    p_c2.font.bold = True
    p_c2.font.color.rgb = COLOR_AZURE

    p_c3 = tf_c.add_paragraph()
    p_c3.text = "GitHub Repository: github.com/Jatinkumar2503/Nexus--AI"
    p_c3.font.size = Pt(12)
    p_c3.font.color.rgb = COLOR_TEXT_SUB

    add_notes(s10, "SPEAKER NOTES:\nThank you judges for your time and consideration. NEXUS AI is ready to pioneer autonomous decision intelligence for railway systems and beyond. I am Jatin Kumar, Team Leader of Team Polaris, and we welcome your questions.")

    # SAVE PRESENTATION TO MULTIPLE PATHS WITH PERMISSION FALLBACK
    output_paths = [
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_OpenLayout_10_Slides.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_WorldClass_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_CODEFUSION_2k26_Polaris.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Best_PPT.pptx"
    ]
    for path in output_paths:
        try:
            prs.save(path)
            print(f"Open Layout 10-Slide PPT successfully saved to: {path}")
        except PermissionError:
            print(f"Skipped locked file: {path} (File is open in PowerPoint)")
        except Exception as e:
            print(f"Error saving to {path}: {e}")

if __name__ == "__main__":
    create_open_layout_10_slide_ppt()
