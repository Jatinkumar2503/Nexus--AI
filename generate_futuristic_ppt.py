import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_futuristic_10_slide_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # HIGH-CONTRAST FUTURISTIC DARK MODE COLOR PALETTE
    COLOR_BG = RGBColor(11, 15, 25)               # Deep Cyber Space #0B0F19
    COLOR_CARD = RGBColor(22, 30, 46)             # Dark Cyber Slate #161E2E
    COLOR_CARD_ALT = RGBColor(30, 41, 59)         # Dark Navy Slate #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85)      # Muted Slate Border #334155
    
    COLOR_TEXT_MAIN = RGBColor(255, 255, 255)     # Pure White #FFFFFF
    COLOR_TEXT_MUTED = RGBColor(203, 213, 225)    # Light Slate Gray #CBD5E1
    COLOR_TEXT_SUB = RGBColor(148, 163, 184)      # Slate Gray #94A3B8
    
    # NEON GLOW ACCENT COLORS
    COLOR_NEON_CYAN = RGBColor(6, 182, 212)       # Neon Cyan #06B6D4
    COLOR_ELECTRIC_BLUE = RGBColor(56, 189, 248)  # Electric Sky Blue #38BDF8
    COLOR_NEON_PURPLE = RGBColor(168, 85, 247)    # Cyber Purple #A855F7
    COLOR_NEON_GREEN = RGBColor(52, 211, 153)     # Glowing Emerald #34D399
    COLOR_NEON_ORANGE = RGBColor(251, 146, 60)    # Bright Coral Orange #FB923C
    COLOR_NEON_RED = RGBColor(248, 113, 113)      # Laser Red #F87171
    COLOR_NEON_YELLOW = RGBColor(250, 204, 21)    # Cyber Gold #FACC15

    def set_dark_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, slide_num):
        # Top Accent Glow Line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = COLOR_NEON_CYAN
        top_line.line.fill.background()

        # Category Tag Badge
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9.0), Inches(0.3))
        tf_tag = tag_box.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f"⚡ {tag_text.upper()}"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_NEON_CYAN

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.5), Inches(0.6))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

        # Counter Badge Box
        c_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.233), Inches(0.35), Inches(1.3), Inches(0.4))
        c_badge.fill.solid()
        c_badge.fill.fore_color.rgb = COLOR_CARD_ALT
        c_badge.line.color.rgb = COLOR_NEON_CYAN
        c_badge.line.width = Pt(1)
        
        tb_b = slide.shapes.add_textbox(Inches(11.233), Inches(0.37), Inches(1.3), Inches(0.35))
        p_b = tb_b.text_frame.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        p_b.text = f"SLIDE {slide_num}/10"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_NEON_CYAN

        # Footer Line & Text
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "NEXUS AI  •  TEAM POLARIS  •  DECISION INTELLIGENCE PLATFORM  •  CODEFUSION 2k26"
        p_foot.font.size = Pt(9)
        p_foot.font.bold = True
        p_foot.font.color.rgb = COLOR_TEXT_SUB

    def add_card(slide, left, top, width, height, border_color=COLOR_CARD_BORDER, bg_color=COLOR_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    def add_notes(slide, notes_text):
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
        except Exception as e:
            pass

    # =========================================================================
    # SLIDE 1: HERO TITLE SLIDE (DARK MODE CYBERPUNK/APPLE TECH)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s1)

    # Hero Background Card
    c_hero = add_card(s1, 0.8, 0.7, 11.733, 6.0, COLOR_ELECTRIC_BLUE, COLOR_CARD)

    # Status Pill Badge
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.1), Inches(4.5), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLOR_CARD_ALT
    pill.line.color.rgb = COLOR_NEON_GREEN
    pill.line.width = Pt(1.2)
    
    tb_pill = s1.shapes.add_textbox(Inches(1.3), Inches(1.15), Inches(4.3), Inches(0.3))
    p = tb_pill.text_frame.paragraphs[0]
    p.text = "🟢 LIVE AI PLATFORM  |  DECISION INTELLIGENCE"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_NEON_GREEN

    # Giant Headline
    tb_t = s1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.5), Inches(1.2))
    p = tb_t.text_frame.paragraphs[0]
    p.text = "NEXUS AI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    # Subtitle
    tb_s = s1.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(10.5), Inches(0.8))
    p = tb_s.text_frame.paragraphs[0]
    p.text = "AI-Native Decision Intelligence & VCG Token Auction Engine"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ELECTRIC_BLUE

    # Tagline
    tb_d = s1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.5), Inches(0.7))
    p = tb_d.text_frame.paragraphs[0]
    p.text = "Autonomously predicting, simulating, and resolving critical railway network bottlenecks before catastrophic operational failure."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # Team Polaris Metadata Card (High Visual Prominence)
    meta_card = add_card(s1, 1.2, 4.4, 10.933, 1.9, COLOR_NEON_CYAN, COLOR_CARD_ALT)
    
    tb_m = s1.shapes.add_textbox(Inches(1.4), Inches(4.55), Inches(10.5), Inches(1.6))
    tf_m = tb_m.text_frame
    
    p1 = tf_m.paragraphs[0]
    p1.text = "🚀 TEAM: POLARIS   |   👑 TEAM LEADER: JATIN KUMAR"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    
    p2 = tf_m.add_paragraph()
    p2.text = "\n🎓 Institution: Deenbandhu Chhotu Ram University of Science & Technology (DCRUST), Murthal"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_NEON_CYAN

    p3 = tf_m.add_paragraph()
    p3.text = "🌐 GitHub: github.com/Jatinkumar2503/Nexus--AI   •   Theme: AI & Infrastructure Resilience"
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s1, "SPEAKER NOTES (Jatin Kumar):\nGood morning judges and attendees! I am Jatin Kumar, leader of Team Polaris from DCRUST Murthal. Today we present NEXUS AI—an AI-native decision intelligence platform designed to protect and optimize critical railway infrastructure. Global power grids and transport networks lose over $100 Billion annually due to delayed, manual dispatcher guesswork during unexpected track disruptions. NEXUS AI solves this by introducing a multi-agent AI swarm, SimPy digital twin, and VCG token slot auctions.")

    # =========================================================================
    # SLIDE 2: THE DISRUPTION CRISIS (HIGH-IMPACT VISUAL)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s2)
    add_header(s2, "01 / THE DISRUPTION CRISIS", "The Cascading Collapse of Manual Infrastructure Operations", 2)

    # Subtitle
    tb_sub = s2.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.733), Inches(0.4))
    p = tb_sub.text_frame.paragraphs[0]
    p.text = "A single 15-minute signal failure triggers a multi-hour domino effect across dense railway corridors."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MUTED

    crisis_cards = [
        ("🔴 1. Signal Glitch", "Track circuit failure at a critical junction node.", COLOR_NEON_RED),
        ("⚠️ 2. Telemetry Overload", "Thousands of raw IoT alerts overwhelm human operators.", COLOR_NEON_ORANGE),
        ("⏳ 3. Manual Guesswork", "Siloed phone calls and manual route recalculation.", COLOR_NEON_YELLOW),
        ("❌ 4. Crew & Safety Breach", "Train crews exceed legal shift limits; gridlock spikes.", COLOR_NEON_RED),
        ("💥 5. $100B+ Global Damage", "Cascading delays burn traction power & destroy passenger trust.", COLOR_NEON_RED)
    ]

    for i, (title, desc, color) in enumerate(crisis_cards):
        c_left = 0.8 + i * 2.38
        c = add_card(s2, c_left, 1.8, 2.18, 4.9, color, COLOR_CARD)
        
        tb = s2.shapes.add_textbox(Inches(c_left + 0.1), Inches(2.0), Inches(1.98), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s2, "SPEAKER NOTES:\nWhen a signal fails or a track is blocked, human dispatchers must manually analyze complex train schedules under intense pressure. This causes delayed reactions, resulting in cascading train standstills, crew overtime violations, and massive energy waste.")

    # =========================================================================
    # SLIDE 3: THE NEXUS AI BREAKTHROUGH
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s3)
    add_header(s3, "02 / THE SOLUTION ARCHITECTURE", "Predict, Simulate & Decide Before Impact", 3)

    pillars = [
        ("⚡ SimPy Physics Digital Twin", "Real-time discrete-event simulation modeling exact train kinematics, platform track capacities, and IEEE catenary power limits.", COLOR_NEON_CYAN),
        ("🎫 VCG Token Slot Auction Engine", "Game-theoretic priority token bidding allocating track slots fairly with second-price pricing and anti-starvation mechanics.", COLOR_NEON_PURPLE),
        ("🤖 LangGraph Multi-Agent Swarm", "Cyclic state machine orchestrating Planner, Risk, Energy, and Safety Validation agents to negotiate Pareto-optimal recovery plans.", COLOR_NEON_GREEN)
    ]

    for i, (title, desc, color) in enumerate(pillars):
        top = 1.75 + i * 1.65
        c = add_card(s3, 0.8, top, 11.733, 1.4, color, COLOR_CARD)
        
        tb = s3.shapes.add_textbox(Inches(1.1), Inches(top + 0.15), Inches(11.1), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s3, "SPEAKER NOTES:\nNEXUS AI introduces three core breakthroughs: A real-time SimPy Digital Twin that models network physics, a VCG Token Auction Engine that allocates track priority mathematically, and a LangGraph Multi-Agent Swarm that generates validated recovery scenarios.")

    # =========================================================================
    # SLIDE 4: SYSTEM PIPELINE FLOWCHART
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s4)
    add_header(s4, "03 / SYSTEM PIPELINE", "End-to-End Edge Telemetry to Cockpit Flowchart", 4)

    pipeline_nodes = [
        ("📡 1. Track Sensors & Azure IoT", "Ingests live telemetry, train velocity, and track blockage events.", COLOR_NEON_CYAN),
        ("🔄 2. SimPy Digital Twin", "Simulates train movement, platform capacity, & power grid load.", COLOR_ELECTRIC_BLUE),
        ("🤖 3. LangGraph AI Swarm", "Multi-agent negotiation formulating Pareto recovery options.", COLOR_NEON_PURPLE),
        ("🎫 4. VCG Auction Solver", "Executes second-price slot bidding & anti-starvation interest.", COLOR_NEON_ORANGE),
        ("🛡️ 5. Dispatcher Cockpit GUI", "Presents validated recommendations for 1-click human approval.", COLOR_NEON_GREEN)
    ]

    for i, (title, desc, color) in enumerate(pipeline_nodes):
        top = 1.6 + i * 1.05
        c = add_card(s4, 0.8, top, 11.733, 0.9, color, COLOR_CARD)
        
        tb = s4.shapes.add_textbox(Inches(1.1), Inches(top + 0.12), Inches(11.1), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s4, "SPEAKER NOTES:\nTelemetry flows from edge track sensors to our SimPy Digital Twin. The LangGraph Multi-Agent Swarm negotiates recovery options, the VCG Solver resolves track slot conflicts, and the optimal strategy is presented to the human dispatcher.")

    # =========================================================================
    # SLIDE 5: VCG TOKEN AUCTION ENGINE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s5)
    add_header(s5, "04 / GAME THEORY ENGINE", "Vickrey-Clarke-Groves (VCG) Token Slot Auctions", 5)

    vcg_cards = [
        ("🎫 Priority Tokens (`tkn`)", "Train agents submit priority tokens reflecting delay impact. High-priority express trains bid higher for critical platform track access.", COLOR_NEON_CYAN),
        ("⚖️ Second-Price Rule", "The highest bidder wins the track slot, but pays only the second-highest bid. Guarantees truthful bidding where no train benefits from over-reporting.", COLOR_NEON_PURPLE),
        ("⏳ Anti-Starvation Rule", "Trains waiting at red signals continuously accrue token interest over time, ensuring low-priority freight trains eventually win slots.", COLOR_NEON_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(vcg_cards):
        left = 0.8 + i * 3.96
        c = add_card(s5, left, 1.8, 3.75, 4.9, color, COLOR_CARD)
        
        tb = s5.shapes.add_textbox(Inches(left + 0.2), Inches(2.0), Inches(3.35), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s5, "SPEAKER NOTES:\nTo solve platform track conflicts fairly, NEXUS AI employs VCG token auctions. Train agents submit priority tokens based on delay penalties. Second-price bidding encourages truthful reporting, while interest token accumulation prevents lower-priority trains from starving.")

    # =========================================================================
    # SLIDE 6: AUTONOMOUS MULTI-AGENT SWARM
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s6)
    add_header(s6, "05 / MULTI-AGENT INTELLIGENCE", "LangGraph Cyclic Multi-Agent Negotiation Swarm", 6)

    agents = [
        ("🗺️ Planner Agent", "Formulates multi-path recovery hypotheses (Detour, Short-Turn, Holds).", COLOR_NEON_CYAN, 0.8, 1.8),
        ("⚡ Risk & Energy Agent", "Calculates traction acceleration power, fuel burn & IEEE catenary limits.", COLOR_ELECTRIC_BLUE, 6.8, 1.8),
        ("🎫 VCG Auction Solver", "Executes game-theoretic second-price priority token bidding.", COLOR_NEON_ORANGE, 0.8, 4.3),
        ("🛡️ Validation & Safety Agent", "Enforces 100% deterministic safety rules & self-reflection bounds.", COLOR_NEON_GREEN, 6.8, 4.3)
    ]

    for title, desc, color, left, top in agents:
        c = add_card(s6, left, top, 5.733, 2.2, color, COLOR_CARD)
        
        tb = s6.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(5.333), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s6, "SPEAKER NOTES:\nOur AI swarm consists of specialized agents working in a cyclic LangGraph state machine. The Planner proposes routes, Risk & Energy verify physical limits, VCG Auction resolves slot order, and Validation guarantees safety before dispatcher review.")

    # =========================================================================
    # SLIDE 7: SCENARIO COMPARISON MATRIX
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s7)
    add_header(s7, "06 / DECISION INTELLIGENCE", "Multi-Objective Pareto-Optimal Tradeoff Analysis", 7)

    options = [
        ("❌ Option A: Do Nothing", "High delay penalty (140+ mins), severe crew overtime violation risk, station gridlock.", COLOR_NEON_RED),
        ("🟢 Option B: Detour Route (RECOMMENDED)", "Pareto Optimal: Moderate extra distance, 0 min station gridlock, optimal safety score.", COLOR_NEON_GREEN),
        ("⚠️ Option C: Short-Turn", "Turns train back early; limits passenger reach but instantly frees platform capacity.", COLOR_NEON_YELLOW)
    ]

    for i, (title, desc, color) in enumerate(options):
        left = 0.8 + i * 3.96
        c = add_card(s7, left, 1.8, 3.75, 4.9, color, COLOR_CARD)
        
        tb = s7.shapes.add_textbox(Inches(left + 0.2), Inches(2.0), Inches(3.35), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s7, "SPEAKER NOTES:\nNEXUS AI doesn't just offer one solution—it provides side-by-side Pareto-optimal comparisons. Here, Option B (Detour Route) is recommended because it completely eliminates station gridlock while maintaining full safety scores.")

    # =========================================================================
    # SLIDE 8: OPERATOR COCKPIT & SAFEGUARDS
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s8)
    add_header(s8, "07 / OPERATOR EXPERIENCE", "Human-in-the-Loop Cockpit & Durable Audit Memory", 8)

    features = [
        ("🗺️ Interactive Spatial Corridor Map", "MapLibre GL visual rendering of live trains, signals, and dynamic track blockages.", COLOR_NEON_CYAN),
        ("🧠 Plain-Language AI Reasoning", "Explains *why* a recovery plan was selected with side-by-side tradeoff analysis.", COLOR_NEON_PURPLE),
        ("🔒 1-Click Dispatcher Approval Gate", "Simulation resumes only after the human operator explicitly approves the recovery action.", COLOR_NEON_GREEN),
        ("📜 Durable Event Replay Memory", "Immutable audit log tracking every incident telemetry event, token bid, and executed action.", COLOR_NEON_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(features):
        top = 1.8 + i * 1.2
        c = add_card(s8, 0.8, top, 11.733, 1.0, color, COLOR_CARD)
        
        tb = s8.shapes.add_textbox(Inches(1.1), Inches(top + 0.12), Inches(11.1), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s8, "SPEAKER NOTES:\nThe React dispatcher cockpit keeps human operators in complete control. It displays spatial train positions, plain-language reasoning, requires explicit dispatcher approval, and maintains an immutable audit trail for post-incident review.")

    # =========================================================================
    # SLIDE 9: EMPIRICAL BENCHMARKS & TECH STACK (318M FOUNDATION MODEL)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s9)
    add_header(s9, "08 / BENCHMARKS & TECH STACK", "318M Foundation Model Validation & Architecture", 9)

    # Hero Metrics Box
    c_m = add_card(s9, 0.8, 1.8, 5.6, 4.9, COLOR_NEON_CYAN, COLOR_CARD)
    tb_m = s9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    p = tf_m.paragraphs[0]
    p.text = "📈 318M FOUNDATION MODEL BENCHMARKS"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_NEON_CYAN
    
    benchmarks = [
        ("🎯 99.65% Dispatch Policy Accuracy", "Trained on 100,000 scenarios; outperforms FIFO heuristics (16.7%) by +51.4%."),
        ("⏱️ 0.1599 min (~9.6s) Delay MAE", "18.5x error reduction over linear tabular models in continuous delay forecasting."),
        ("⚡ 1.39 ms TorchScript JIT Inference", "Optimized C++ runtime for real-time edge control inside interlocking stations."),
        ("📉 34.5% Historical Delay Reduction", "Backtested on Northern Fog Gridlock & Western Monsoon flooding incidents.")
    ]
    for title, desc in benchmarks:
        p1 = tf_m.add_paragraph()
        p1.text = f"\n{title}"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p2 = tf_m.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Tech Stack Box
    c_t = add_card(s9, 6.933, 1.8, 5.6, 4.9, COLOR_NEON_PURPLE, COLOR_CARD)
    tb_t = s9.shapes.add_textbox(Inches(7.133), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    
    p = tf_t.paragraphs[0]
    p.text = "💻 PRODUCTION TECH STACK"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_NEON_PURPLE

    stack = [
        ("🧠 Deep Learning Core", "PyTorch 2.0+, Hetero-GAT, Causal Temporal Transformers (318M params)"),
        ("⚡ Optimization & Serving", "TorchScript JIT (1.39ms P50), INT8 Quantization, FastAPI, WebSockets"),
        ("🎮 Digital Twin & Graph", "SimPy Physics Simulator, NetworkX Corridor Graphs, MapLibre GL"),
        ("🎨 Cockpit UI & Cloud", "React 18, TypeScript, Tailwind CSS, Docker, Azure IoT Hub")
    ]
    for title, desc in stack:
        p1 = tf_t.add_paragraph()
        p1.text = f"\n{title}"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p2 = tf_t.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(s9, "SPEAKER NOTES:\nOur 318M Spatiotemporal Foundation Model achieves 99.65% top-1 dispatch accuracy and sub-10 second delay MAE across 100,000 scenarios. Compiled with TorchScript JIT, it executes in just 1.39 ms P50 latency, delivering 34.5% delay reduction over manual human dispatch.")

    # =========================================================================
    # SLIDE 10: STRATEGIC ADVANTAGE & THANK YOU
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_dark_bg(s10)
    add_header(s10, "09 / SUMMARY & CONCLUSION", "Transforming Critical Infrastructure Operations", 10)

    # Hero Card
    c_hero = add_card(s10, 0.8, 1.8, 11.733, 4.9, COLOR_NEON_CYAN, COLOR_CARD)
    
    tb_h = s10.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.933), Inches(4.4))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    
    p = tf_h.paragraphs[0]
    p.text = "🏆 WHY NEXUS AI OUTPERFORMS COMPETITORS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_NEON_CYAN

    points = [
        "1. 🤖 AI-Native Architecture: Purpose-built multi-agent swarm intelligence for critical infrastructure.",
        "2. ⚖️ Fair VCG Auctions: Zero starvation, truthful slot allocation for competing train priorities.",
        "3. 🛡️ Deterministic Safety Safeguards: Local fallback engine guarantees 100% operational compliance.",
        "4. ⚡ Sub-Second Simulation Speed: SimPy digital twin processes 10,000+ complex nodes instantaneously."
    ]
    for pt in points:
        p_pt = tf_h.add_paragraph()
        p_pt.text = f"\n{pt}"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = COLOR_TEXT_MUTED

    # Contact Card
    contact_box = add_card(s10, 1.2, 4.9, 10.933, 1.4, COLOR_NEON_GREEN, COLOR_CARD_ALT)
    tb_c = s10.shapes.add_textbox(Inches(1.4), Inches(5.05), Inches(10.5), Inches(1.2))
    tf_c = tb_c.text_frame
    
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "🚀 THANK YOU  |  TEAM: POLARIS  |  👑 TEAM LEADER: JATIN KUMAR"
    p_c1.font.size = Pt(16)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_TEXT_MAIN

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "🎓 Institution: DCRUST Murthal   •   🌐 GitHub Repository: github.com/Jatinkumar2503/Nexus--AI"
    p_c2.font.size = Pt(12)
    p_c2.font.bold = True
    p_c2.font.color.rgb = COLOR_NEON_CYAN

    add_notes(s10, "SPEAKER NOTES:\nThank you judges for your time and consideration. NEXUS AI is ready to pioneer autonomous decision intelligence for railway systems and beyond. I am Jatin Kumar, Team Leader of Team Polaris, and we welcome your questions.")

    # SAVE PRESENTATION TO MULTIPLE PATHS WITH PERMISSION FALLBACK
    output_paths = [
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Futuristic_Deck.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_10_Slides_Futuristic.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_10_Slides.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_WorldClass_Presentation.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_CODEFUSION_2k26_Polaris.pptx",
        r"c:\Users\Asus\Documents\far away\NEXUS_AI_Polaris_Best_PPT.pptx"
    ]
    for path in output_paths:
        try:
            prs.save(path)
            print(f"Futuristic Dark-Mode 10-Slide PPT successfully saved to: {path}")
        except PermissionError:
            print(f"Skipped locked file: {path} (File is currently open in PowerPoint)")
        except Exception as e:
            print(f"Error saving to {path}: {e}")

if __name__ == "__main__":
    create_futuristic_10_slide_ppt()
